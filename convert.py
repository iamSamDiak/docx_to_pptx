from docx import Document
from docx.opc.exceptions import PackageNotFoundError
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import re

class Powerpoint:
    def __init__(self):
        self.pptx = Presentation()
        self.paragraphs = []

    def open(self, path_of_docx):
        """ Ouvre un document Word """
        self.document = None
        try:
            self.path_of_docx = path_of_docx
            self.document = Document(self.path_of_docx)
            print("Fichier trouvé")
        except PackageNotFoundError:
            print("Fichier introuvable!")
        except:
            print("Erreur inconnue! Contactez-moi.")

    def get_text(self):
        """ Récupère tous les textes [ [{ref, text}] ] """
        arr = []
        p = { "headline": "", "texts": "" }
        if self.document == None:
            return
        for i in self.document.paragraphs:
            size = [run.font.size / 12700 if run.font.size else None for run in i.runs]
            is_headline = any(val is not None and val >= 16.0 for val in size)

            """ Si le texte est un titre, crée un nouveau 'dict' """
            if is_headline:
                if p["headline"]:
                    arr.append(p)
                    p = { "headline": "", "texts": "" }
                p["headline"] = i.text
            else:
                if i.text != "" and i.text is not None:
                    if p["texts"]:
                        p["texts"] = p["texts"] + "\n" + i.text
                    else:
                        p["texts"] = p["texts"] + i.text

        if p["headline"] or p["texts"]:
            arr.append(p)
        return arr

    def add_textbox(self, top):
        slide_width = self.pptx.slide_width
        slide_height = self.pptx.slide_height
        textbox_width = slide_width - 50 * 12700
        left = (slide_width - textbox_width) / 2
        txBox = self.slide.shapes.add_textbox(left, Pt(top), textbox_width, slide_height)
        tf = txBox.text_frame
        tf.word_wrap = True
        return tf

    def to_pptx(self):
        """ Crée des slides par rapport aux textes """
        print("Conversion...")
        self.paragraphs = self.get_text()
        #
        for para in self.paragraphs:
            headline = para["headline"]
            # chunk
            chunk = self.smart_split_text(para["texts"], 520)
            for index, text in enumerate(chunk):
                blank_slide_layout = self.pptx.slide_layouts[6]
                self.slide = self.pptx.slides.add_slide(blank_slide_layout)
                # headline
                textbox_headline = self.add_textbox(30)
                textbox_headline.text = headline + f"\n({ index + 1 }/{ len(chunk) })" if len(chunk) > 1 else headline
                for p in textbox_headline.paragraphs:
                    p.alignment = PP_ALIGN.CENTER
                    run = p.runs[0]
                    run.font.bold = True
                    run.font.size = Pt(26)
                    run.font.color.rgb = RGBColor(33, 55, 87)
                # text
                textbox_chunk = self.add_textbox(125) if len(chunk) > 1 else self.add_textbox(110)
                textbox_chunk.text = text
                for p in textbox_chunk.paragraphs:
                    p.alignment = PP_ALIGN.JUSTIFY
                    p.line_spacing = 1.5
                    run = p.runs[0]
                    run.font.bold = True
                    run.font.size = Pt(22)
                    run.font.color.rgb = RGBColor(0, 0, 0)
        print("Conversion finie")
    
    def smart_split_text(self, text, limit):
        """" Si le texte est trop long, il est divisé en plusieurs parties """
        arr_text = []
        start = 0
        total = len(text)
        stop = limit
        if total <= limit:
            arr_text.append(text.lstrip("\n").lstrip(" "))
        else:
            while start < total:
                end = min(start + limit, total)
                # Cherche la prochaine ponctuation après la limite
                match = re.search(r'[.,;:!?]', text[end:])
                if match:
                    split_at = end + match.start() + 1
                    arr_text.append(text[start:split_at].lstrip("\n").lstrip(" "))
                    start = split_at
                else:
                    arr_text.append(text[start:].lstrip("\n").lstrip(" "))
                    break
        return arr_text

    def save(self):
        """ Exporte le fichier Powerpoint """
        self.pptx.save("Test2.pptx")
        print("Powerpoint exporté")

pw = Powerpoint()
pw.open("Un test.docx")
pw.get_text()
pw.to_pptx()
pw.save()